"""Create a PowerPoint PPTX file from JSON slide data using python-pptx."""

import json
import os


def run(file_path: str, slides: str) -> str:
    """
    Create a PowerPoint file from a JSON array of slide objects.

    Each slide object should have a 'title' (str) and 'content' (str) field.

    @param file_path: Output path for the PPTX file.
    @param slides: JSON string, e.g. '[{"title":"Intro","content":"Welcome text"}]'.
    @returns: Confirmation message with the output path.
    @throws ImportError: If python-pptx is not installed.
    @throws ValueError: If slides is not valid JSON or has incorrect structure.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "error: " + "python-pptx is required but not installed. "
            "Install it with: pip install python-pptx"

    try:
        slide_data = json.loads(slides)
    except json.JSONDecodeError as err:
        raise ValueError(f"Invalid JSON slides data: {err}")

    if not isinstance(slide_data, list):
        raise ValueError("Slides must be a JSON array of slide objects.")

    resolved = os.path.abspath(file_path)
    parent_dir = os.path.dirname(resolved)
    if parent_dir and not os.path.exists(parent_dir):
        os.makedirs(parent_dir, exist_ok=True)

    presentation = Presentation()

    for index, slide_info in enumerate(slide_data):
        if not isinstance(slide_info, dict):
            raise ValueError(
                f"Slide at index {index} must be an object with 'title' and 'content'."
            )
        title = slide_info.get("title", "")
        content = slide_info.get("content", "")

        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 1:
                placeholder.text = content
                break

    presentation.save(resolved)
    return f"PowerPoint file created successfully at: {resolved}"
