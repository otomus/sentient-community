"""Read text content from a PowerPoint PPTX file using python-pptx."""

import os


def run(file_path: str) -> str:
    """
    Extract all text content from a PowerPoint file.

    @param file_path: Path to the PPTX file to read.
    @returns: Extracted text organized by slide number.
    @throws FileNotFoundError: If the PPTX file does not exist.
    @throws ImportError: If python-pptx is not installed.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "error: " + "python-pptx is required but not installed. "
            "Install it with: pip install python-pptx"

    resolved = os.path.abspath(file_path)
    if not os.path.exists(resolved):
        raise FileNotFoundError(f"PowerPoint file not found: {resolved}")

    presentation = Presentation(resolved)
    slides_text: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = _extract_slide_text(slide)
        if texts:
            slide_content = "\n".join(texts)
            slides_text.append(f"--- Slide {slide_number} ---\n{slide_content}")

    if not slides_text:
        return "No text content found in the PowerPoint file."

    return "\n\n".join(slides_text)


def _extract_slide_text(slide: object) -> list[str]:
    """Extract all text fragments from a single slide's shapes."""
    texts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)
    return texts
